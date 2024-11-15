import json
import os
from glob import glob

import jsonlines
import Levenshtein
from jinja2 import Template
from PIL import Image
from pptx import Presentation
from torch import cosine_similarity
from transformers import CLIPModel, CLIPProcessor

import llms

outline_template = Template(
    """
    From the following text which contains a set of headings and some content within each heading:
    {{ text }}
    Extract the most important headings present in it. Reduce the length of each heading to five words if they are lengthy.
    Example Output:
    ["Heading 1", "Heading 2", "Heading 3"]
    Output:

    """
)
mapping_template = Template(
    """
    Think step by step and then answer the following question:  You are given with the following title: {{outline_headings}}
    and a list of keys: {{document_heading_from_bird_eye_view}}
    Each key is associated with some text as presented in the dictionary format below:
    {{bird_eye_view}}
    The task is to find 1-2 significantly matched keys. The matching should be done based on the similarity of the text associated with the keys with the given heading.
    Example Output:
    {"Heading 1": ["key1", "key2"], "Heading 2": ["key1", "key4"]}
    Output:
    """
)
generation_template = Template(
    """
    You are a presentation generator from a source of text. You have to generate the slide number {{slide_index}}. Previous slide headings and slide contents are given below in the format of a list of dictionaries. {{previous_slide}} Given the following slide heading and the source of text respectively, create the content of the slide number {{slide_index}} such that: 1. The slide should have maximum {{max_bullet}} bullet points. 2. Ensure that the content of the bullet points are coming strictly from the given source of text only. 3. The content of the slide is very relevant to the given slide heading 4. Each bullet point should have a maximum of 10 words 5. Ensure that this slide does not have any content repeated from the previous slides. 6. The flow of the overall presentation is nice. 7. Do not prefix the slide title before the bullet poide nts in the output  SliTitle: {{slide_heading}} Source of text: {{text}}
    Example Output:
    ["bullet point 1", "bullet point 2"]
    Output: give your output as a list of strings in json format
    """
)


def filter_aspect_ratio(image: list[str]):
    filtered_images = []
    for i in image:
        size = image = Image.open(i).size
        long, short = max(size), min(size)
        if long / short < 4:
            filtered_images.append(i)
    return filtered_images


def edit_distance(text1: str, text2: str):
    return 1 - Levenshtein.distance(text1, text2) / max(len(text1), len(text2))


def generate_content(source_text: str, bird_eye: dict, max_bullet: int):
    bird_eye_headdings = []
    for section in bird_eye["sections"]:
        bird_eye_headdings.extend(
            [next(iter(subsec)) for subsec in section["subsections"]]
        )
    outline: list[str] = llms.language_model(
        outline_template.render(text=source_text), return_json=True
    )
    assert len(outline) != 0, "No outline found"
    mapping = llms.language_model(
        mapping_template.render(
            outline_headings=outline,
            document_heading_from_bird_eye_view=bird_eye_headdings,
            bird_eye_view=bird_eye,
        ),
        return_json=True,
    )
    assert len(mapping) == len(outline), "Mapping not found"
    slides = []
    for slide_title in outline:
        indexed_sections = []
        for section in bird_eye["sections"]:
            for subsection in section["subsections"]:
                if any(
                    edit_distance(key, next(iter(subsection))) > 0.9
                    for key in mapping[slide_title]
                ):
                    indexed_sections.append(subsection)
        bullet_points = llms.language_model(
            generation_template.render(
                slide_heading=slide_title,
                text=indexed_sections,
                previous_slide=slides,
                max_bullet=max_bullet,
            ),
            return_json=True,
        )
        assert len(bullet_points) != 0, f"No bullet points found for {slide_title}"
        slides.append(
            {
                "title": slide_title,
                "bullets": bullet_points,
                "indexed_sections": mapping[slide_title],
            }
        )
    return slides


def generate_slides(
    output_file: str,
    source_text: str,
    bird_eye: dict,
    images: list[str],
    model: CLIPModel,
    processor: CLIPProcessor,
):
    if os.path.exists(output_file + ".json"):
        return
    images = filter_aspect_ratio(images)
    slides = generate_content(source_text, bird_eye, 7)
    image_embeddings = model.get_image_features(
        **processor(images=[Image.open(i) for i in images], return_tensors="pt").to(
            "cuda"
        )
    ).unsqueeze(0)
    text_embeddings = model.get_text_features(
        **processor(
            text=["\n".join(slide["bullets"]) for slide in slides],
            return_tensors="pt",
            padding=True,
            max_length=77,
            truncation=True,
        ).to("cuda")
    ).unsqueeze(1)
    similarity = cosine_similarity(image_embeddings, text_embeddings, dim=-1)
    pptx = Presentation()
    for slide_idx, slide in enumerate(slides):  # match image here
        title = slide["title"]
        bullets = slide["bullets"]

        subsimilarity = similarity[slide_idx]
        if subsimilarity.max() > 0.8:
            slide = pptx.slides.add_slide(pptx.slide_layouts[6])
            bullets_placeholder = slide.shapes.placeholders[2]
            image = images[subsimilarity.argmax()]
            slides[slide_idx]["image"] = image
            slide.shapes.placeholders[1].insert_picture(image)
        else:
            slide = pptx.slides.add_slide(pptx.slide_layouts[1])
            bullets_placeholder = slide.shapes.placeholders[1]
        slide.shapes.title.text = title
        text_frame = bullets_placeholder.text_frame
        for bullet in bullets:
            para = text_frame.add_paragraph()
            para.text = bullet
            para.level = 1
    with jsonlines.open(output_file + ".jsonl", "w") as writer:
        writer.write_all(slides)
    pptx.save(output_file + ".pptx")


if __name__ == "__main__":
    from tqdm.auto import tqdm

    print("Generating slides on baseline with ", llms.language_model.model)
    model = CLIPModel.from_pretrained("openai/clip-vit-large-patch14").to("cuda").eval()
    processor = CLIPProcessor.from_pretrained("openai/clip-vit-large-patch14")
    for pdf_folder in tqdm(glob("data/*/pdf/*")):
        source_text = open(f"{pdf_folder}/source.md").read()
        bird_eye = json.load(open(f"{pdf_folder}/refined_doc.json"))
        images = json.load(open(f"{pdf_folder}/image_caption.json")).keys()
        output_file = f"{pdf_folder}/baseline_{llms.language_model.model.split('-')[0]}"
        try:
            generate_slides(
                output_file,
                source_text,
                bird_eye,
                list(images),
                model,
                processor,
            )
        except Exception as e:
            print(f"Error in {pdf_folder}: {e}")
