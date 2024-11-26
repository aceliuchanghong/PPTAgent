import os
import shutil
import subprocess
import tempfile
import traceback
from time import time
from types import SimpleNamespace

import json_repair
import Levenshtein
from lxml import etree
from pdf2image import convert_from_path
from pptx.dml.color import RGBColor
from pptx.dml.fill import _NoFill, _NoneFill
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.text.text import _Paragraph, _Run
from pptx.util import Length, Pt
from rich import print
from tenacity import RetryCallState, retry, stop_after_attempt, wait_fixed

IMAGE_EXTENSIONS = {"bmp", "jpg", "jpeg", "pgm", "png", "ppm", "tif", "tiff", "webp"}

BLACK = RGBColor(0, 0, 0)
YELLOW = RGBColor(255, 255, 0)
BORDER_LEN = Pt(2)
LABEL_LEN = Pt(24)
FONT_LEN = Pt(20)


def prepare_shape_label(shape_idx: int, shape: BaseShape):
    shape.line.color.rgb = BLACK
    shape.line.width = BORDER_LEN
    left = shape.left - BORDER_LEN
    top = shape.top + shape.height + BORDER_LEN - LABEL_LEN
    textbox = shape._parent.add_textbox(left, top, LABEL_LEN, LABEL_LEN)
    textbox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    textbox.fill.solid()
    textbox.fill.fore_color.rgb = YELLOW

    p = textbox.text_frame.paragraphs[0]
    p.text = str(shape_idx)
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_LEN
    p.font.bold = True


def get_font_style(font: dict):
    font = SimpleNamespace(**font)
    styles = []
    if font.size:
        styles.append(f"font-size: {font.size}pt")
    if font.color:
        styles.append(f"color: #{font.color}")
    if font.bold:
        styles.append("font-weight: bold")
    if font.italic:
        styles.append("font-style: italic")
    return "; ".join(styles)


def runs_merge(paragraph: _Paragraph):
    runs = paragraph.runs
    if len(runs) == 0:
        runs = [
            _Run(r, paragraph)
            for r in parse_xml(paragraph._element.xml.replace("fld", "r")).r_lst
        ]
    if len(runs) < 1:
        return runs
    pre_run = runs[0]
    new_runs = [pre_run]
    pre_font = runs[0].font
    for run in runs[1:]:
        if run.font != pre_font:
            new_runs.append(run)
            pre_run = run
            pre_font = run.font
        else:
            pre_run.text += run.text
            run._r.getparent().remove(run._r)
    return new_runs


def older_than(filepath, seconds: int = 10):
    if not os.path.exists(filepath):
        return False
    file_creation_time = os.path.getctime(filepath)
    current_time = time()
    return seconds < (current_time - file_creation_time)


def edit_distance(text1: str, text2: str):
    return 1 - Levenshtein.distance(text1, text2) / max(len(text1), len(text2))


def get_slide_content(doc_json: dict, slide_title: str, slide: dict):
    slide_desc = slide.get("description", "")
    slide_content = f"Title: {slide_title}\nSlide Description: {slide_desc}\n"
    for key in slide.get("subsection_keys", []):
        slide_content += "Slide Reference: "
        for section in doc_json["sections"]:
            for subsection in section.get("subsections", []):
                if edit_distance(key, subsection["title"]) > 0.9:
                    slide_content += f"# {key} \n{subsection['content']}\n"
    return slide_content


def tenacity_log(retry_state: RetryCallState):
    print(retry_state)
    traceback.print_tb(retry_state.outcome.exception().__traceback__)


def get_json_from_response(raw_response: str):
    response = raw_response.strip()
    l, r = response.rfind("```json"), response.rfind("```")
    try:
        if l == -1 or r == -1:
            response = json_repair.loads(response)
        else:
            response = json_repair.loads(response[l + 7 : r].strip())
        return response
    except Exception as e:
        raise RuntimeError("Failed to parse JSON from response", e)


tenacity = retry(
    wait=wait_fixed(3), stop=stop_after_attempt(3), after=tenacity_log, reraise=True
)


@tenacity
def ppt_to_images(file: str, output_dir: str):
    os.makedirs(output_dir, exist_ok=True)
    with tempfile.TemporaryDirectory() as temp_dir:
        command_list = [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            file,
            "--outdir",
            temp_dir,
        ]
        subprocess.run(command_list, check=True, stdout=subprocess.DEVNULL)

        for f in os.listdir(temp_dir):
            if not f.endswith(".pdf"):
                continue
            temp_pdf = pjoin(temp_dir, f)
            images = convert_from_path(temp_pdf, dpi=72)
            for i, img in enumerate(images):
                img.save(pjoin(output_dir, f"slide_{i+1:04d}.jpg"))
            return

        raise RuntimeError("No PDF file was created in the temporary directory")


def extract_fill(shape: BaseShape):
    if "fill" not in dir(shape):
        return None
    fill_dict = {
        "fill_xml": shape.fill._xPr.xml,
    } | {k: v for k, v in object_to_dict(shape.fill).items() if "color" in k}
    if not isinstance(shape.fill._fill, (_NoneFill, _NoFill)):
        fill_dict["type"] = shape.fill.type.name.lower()
    return fill_dict


def apply_fill(shape: BaseShape, fill: dict):
    if fill is None:
        return
    new_element = etree.fromstring(fill["fill_xml"])
    shape.fill._xPr.getparent().replace(shape.fill._xPr, new_element)


def parse_groupshape(groupshape: GroupShape):
    assert isinstance(groupshape, GroupShape)
    group_top_left_x = groupshape.left
    group_top_left_y = groupshape.top
    group_width = groupshape.width
    group_height = groupshape.height
    shape_top_left_x = min([sp.left for sp in groupshape.shapes])
    shape_top_left_y = min([sp.top for sp in groupshape.shapes])
    shape_width = (
        max([sp.left + sp.width for sp in groupshape.shapes]) - shape_top_left_x
    )
    shape_height = (
        max([sp.top + sp.height for sp in groupshape.shapes]) - shape_top_left_y
    )
    group_shape_xy = []
    for sp in groupshape.shapes:
        group_shape_left = (
            sp.left - shape_top_left_x
        ) * group_width / shape_width + group_top_left_x
        group_shape_top = (
            sp.top - shape_top_left_y
        ) * group_height / shape_height + group_top_left_y
        group_shape_width = sp.width * group_width / shape_width
        group_shape_height = sp.height * group_height / shape_height
        group_shape_xy.append(
            {
                "left": Length(group_shape_left),
                "top": Length(group_shape_top),
                "width": Length(group_shape_width),
                "height": Length(group_shape_height),
            }
        )
    return group_shape_xy


def is_primitive(obj):
    """
    判断对象或该集合包含的所有对象是否是基本类型。

    参数:
    obj: 要判断的对象

    返回:
    如果对象是基本类型，返回True，否则返回False
    """
    if isinstance(obj, (list, tuple, set, frozenset)):
        return all(is_primitive(item) for item in obj)
    return isinstance(
        obj, (int, float, complex, bool, str, bytes, bytearray, type(None))
    )


DEFAULT_EXCLUDE = set(["element", "language_id", "ln", "placeholder_format"])


def object_to_dict(obj, result=None, exclude=None):
    """
    将对象的非隐藏属性拷贝到一个字典中。

    参数:
    obj: 要拷贝属性的对象

    返回:
    包含对象非隐藏属性的字典
    """
    if result is None:
        result = {}
    exclude = DEFAULT_EXCLUDE.union(exclude or set())
    for attr in dir(obj):
        if attr in exclude:
            continue
        try:
            if not attr.startswith("_") and not callable(getattr(obj, attr)):
                attr_value = getattr(obj, attr)
                if "real" in dir(attr_value):
                    attr_value = attr_value.real
                if attr == "size" and isinstance(attr_value, int):
                    attr_value = Length(attr_value).pt

                if is_primitive(attr_value):
                    result[attr] = attr_value
        except:
            pass
    return result


def merge_dict(d1: dict, d2: list[dict]):
    if len(d2) == 0:
        return d1
    for key in list(d1.keys()):
        values = [d[key] for d in d2]
        if d1[key] is not None and len(values) != 1:
            values.append(d1[key])
        if values[0] is None or not all(value == values[0] for value in values):
            continue
        d1[key] = values[0]
        for d in d2:
            d[key] = None
    return d1


def dict_to_object(dict: dict, obj: object, exclude=None):
    """
    从字典中恢复对象的属性。

    参数:
    d: 包含对象属性的字典
    obj: 要恢复属性的对象

    返回:
    恢复属性后的对象
    """
    if exclude is None:
        exclude = set()
    for key, value in dict.items():
        if key not in exclude:
            setattr(obj, key, value)


class Config:

    def __init__(self, rundir=None, session_id=None, debug=True):
        self.DEBUG = debug
        if session_id is not None:
            self.set_session(session_id)
        if rundir is not None:
            self.set_rundir(rundir)

    def set_session(self, session_id):
        self.session_id = session_id
        self.set_rundir(f"./runs/{session_id}")

    def set_rundir(self, rundir: str):
        self.RUN_DIR = rundir
        self.IMAGE_DIR = pjoin(self.RUN_DIR, "images")
        for the_dir in [self.RUN_DIR, self.IMAGE_DIR]:
            os.makedirs(the_dir, exist_ok=True)

    def set_debug(self, debug: bool):
        self.DEBUG = debug

    def remove_rundir(self):
        if pexists(self.RUN_DIR):
            shutil.rmtree(self.RUN_DIR)
        if pexists(self.IMAGE_DIR):
            shutil.rmtree(self.IMAGE_DIR)


pjoin = os.path.join
pexists = os.path.exists
pbasename = os.path.basename

if __name__ == "__main__":
    config = Config()
    print(config)
